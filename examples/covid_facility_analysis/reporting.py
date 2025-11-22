"""
Dashboard and Report Generation for COVID-19 Facility Data.

This module generates:
- Power BI-style dashboards (HTML/interactive)
- PowerPoint executive decks
- Excel workbooks with formatted tables and charts
- PDF summary reports
"""

import pandas as pd
import numpy as np
from datetime import datetime
from typing import Dict, List, Optional
import logging
from pathlib import Path

# Visualization libraries
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
import seaborn as sns
from matplotlib.patches import Rectangle

# PowerPoint generation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# Excel generation
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.chart import LineChart, BarChart, PieChart, Reference

from analysis import CovidFacilityAnalyzer

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Set visualization style
sns.set_style("whitegrid")
plt.rcParams['figure.figsize'] = (12, 6)
plt.rcParams['font.size'] = 10


class DashboardGenerator:
    """
    Generate dashboards and reports for VA COVID-19 response.

    Creates executive dashboards in multiple formats:
    - Interactive HTML dashboards
    - PowerPoint presentations
    - Excel workbooks
    - PDF reports
    """

    def __init__(
        self,
        analyzer: CovidFacilityAnalyzer,
        output_dir: str = "output"
    ):
        """
        Initialize dashboard generator.

        Args:
            analyzer: CovidFacilityAnalyzer instance with data
            output_dir: Directory for output files
        """
        self.analyzer = analyzer
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

        # Generate summary for reuse
        self.summary = analyzer.generate_executive_summary()

        logger.info(f"Dashboard generator initialized. Output: {self.output_dir}")

    # ========================================================================
    # VISUALIZATION CHARTS
    # ========================================================================

    def create_time_series_chart(
        self,
        metric: str,
        title: str,
        ylabel: str,
        save_path: Optional[str] = None
    ) -> str:
        """
        Create time series line chart.

        Args:
            metric: Metric column name
            title: Chart title
            ylabel: Y-axis label
            save_path: Path to save chart (optional)

        Returns:
            Path to saved chart
        """
        logger.info(f"Creating time series chart for '{metric}'")

        # Get time series data
        trends = self.analyzer.analyze_time_trends(metric, window_days=7)

        fig, ax = plt.subplots(figsize=(12, 6))

        # Plot daily values
        ax.plot(
            trends['date'],
            trends[f'{metric}_total'],
            marker='o',
            markersize=4,
            linewidth=1,
            alpha=0.6,
            label='Daily Values'
        )

        # Plot 7-day rolling average
        ax.plot(
            trends['date'],
            trends[f'{metric}_rolling_avg'],
            linewidth=2.5,
            label='7-Day Moving Average',
            color='red'
        )

        ax.set_title(title, fontsize=14, fontweight='bold')
        ax.set_xlabel('Date', fontsize=12)
        ax.set_ylabel(ylabel, fontsize=12)
        ax.legend(loc='best')
        ax.grid(True, alpha=0.3)

        # Format x-axis dates
        ax.xaxis.set_major_formatter(mdates.DateFormatter('%m/%d'))
        ax.xaxis.set_major_locator(mdates.DayLocator(interval=7))
        plt.xticks(rotation=45)

        plt.tight_layout()

        # Save chart
        if save_path is None:
            save_path = self.output_dir / f"chart_timeseries_{metric}.png"

        plt.savefig(save_path, dpi=300, bbox_inches='tight')
        plt.close()

        logger.info(f"Time series chart saved: {save_path}")
        return str(save_path)

    def create_facility_type_comparison_chart(
        self,
        save_path: Optional[str] = None
    ) -> str:
        """
        Create bar chart comparing metrics by facility type.

        Args:
            save_path: Path to save chart

        Returns:
            Path to saved chart
        """
        logger.info("Creating facility type comparison chart")

        comparison = self.analyzer.compare_by_facility_type()

        fig, axes = plt.subplots(2, 2, figsize=(14, 10))
        fig.suptitle('VA Facility Type Comparison - COVID-19 Metrics',
                    fontsize=16, fontweight='bold')

        # Chart 1: COVID Positive by facility type
        ax1 = axes[0, 0]
        ax1.bar(comparison['facility_type'], comparison['covid_positive_total'])
        ax1.set_title('Total COVID-19 Positive Cases')
        ax1.set_ylabel('Cases')
        ax1.tick_params(axis='x', rotation=45)

        # Chart 2: Average occupancy rate
        ax2 = axes[0, 1]
        colors = ['red' if x >= 0.85 else 'orange' if x >= 0.75 else 'green'
                 for x in comparison['avg_occupancy_rate']]
        ax2.bar(comparison['facility_type'], comparison['avg_occupancy_rate'], color=colors)
        ax2.axhline(y=0.85, color='red', linestyle='--', linewidth=2, label='Critical (85%)')
        ax2.axhline(y=0.75, color='orange', linestyle='--', linewidth=2, label='High (75%)')
        ax2.set_title('Average Occupancy Rate')
        ax2.set_ylabel('Occupancy Rate')
        ax2.legend()
        ax2.tick_params(axis='x', rotation=45)

        # Chart 3: COVID Deaths
        ax3 = axes[1, 0]
        ax3.bar(comparison['facility_type'], comparison['covid_deaths_total'])
        ax3.set_title('Total COVID-19 Deaths')
        ax3.set_ylabel('Deaths')
        ax3.tick_params(axis='x', rotation=45)

        # Chart 4: Facilities with constraints
        ax4 = axes[1, 1]
        x = np.arange(len(comparison))
        width = 0.35
        ax4.bar(x - width/2, comparison['facilities_with_staff_shortage'],
               width, label='Staff Shortage')
        ax4.bar(x + width/2, comparison['facilities_with_ppe_critical'],
               width, label='PPE Critical')
        ax4.set_title('Facilities with Resource Constraints')
        ax4.set_ylabel('Number of Facilities')
        ax4.set_xticks(x)
        ax4.set_xticklabels(comparison['facility_type'], rotation=45)
        ax4.legend()

        plt.tight_layout()

        if save_path is None:
            save_path = self.output_dir / "chart_facility_type_comparison.png"

        plt.savefig(save_path, dpi=300, bbox_inches='tight')
        plt.close()

        logger.info(f"Facility type comparison chart saved: {save_path}")
        return str(save_path)

    def create_kpi_dashboard_chart(
        self,
        save_path: Optional[str] = None
    ) -> str:
        """
        Create KPI dashboard with key metrics.

        Args:
            save_path: Path to save chart

        Returns:
            Path to saved chart
        """
        logger.info("Creating KPI dashboard chart")

        kpis = self.summary['kpis']

        fig = plt.figure(figsize=(16, 10))
        gs = fig.add_gridspec(3, 4, hspace=0.4, wspace=0.3)

        # Title
        fig.suptitle(
            f'VA COVID-19 Executive Dashboard - {kpis["report_date"]}',
            fontsize=18,
            fontweight='bold',
            y=0.98
        )

        # KPI boxes
        kpi_data = [
            ("COVID-19 Positive", kpis['total_covid_positive'], "cases", "#FF6B6B"),
            ("Hospitalized", kpis['total_covid_hospitalized'], "patients", "#4ECDC4"),
            ("ICU", kpis['total_covid_icu'], "patients", "#95E1D3"),
            ("Ventilators", kpis['total_covid_ventilator'], "patients", "#F38181"),
            ("Deaths", kpis['total_covid_deaths'], "deaths", "#AA4465"),
            ("Tests Conducted", kpis['total_tests_conducted'], "tests", "#A8E6CF"),
            ("Positivity Rate", f"{kpis['avg_positivity_rate']:.1%}", "", "#FFD93D"),
            ("Avg Occupancy", f"{kpis['avg_occupancy_rate']:.1%}", "", "#6BCB77"),
        ]

        for idx, (label, value, unit, color) in enumerate(kpi_data[:8]):
            row = idx // 4
            col = idx % 4
            ax = fig.add_subplot(gs[row, col])
            ax.axis('off')

            # Create KPI box
            rect = Rectangle((0, 0), 1, 1, facecolor=color, alpha=0.7, edgecolor='black', linewidth=2)
            ax.add_patch(rect)

            # Add text
            if isinstance(value, int):
                value_str = f"{value:,}"
            else:
                value_str = str(value)

            ax.text(0.5, 0.65, value_str, ha='center', va='center',
                   fontsize=24, fontweight='bold', color='white')
            ax.text(0.5, 0.35, f"{label}\n{unit}", ha='center', va='center',
                   fontsize=11, color='white', style='italic')

            ax.set_xlim(0, 1)
            ax.set_ylim(0, 1)

        # Bottom row: Resource constraints
        ax_bottom = fig.add_subplot(gs[2, :])

        constraints = [
            f"Staff Shortages: {kpis['facilities_staff_shortage']} facilities",
            f"PPE Critical: {kpis['facilities_ppe_critical']} facilities",
            f"Avg N95 Supply: {kpis['avg_n95_days_supply']:.1f} days",
            f"Data Completeness: {kpis['avg_data_completeness']:.1%}",
        ]

        constraint_text = " | ".join(constraints)
        ax_bottom.text(0.5, 0.5, constraint_text, ha='center', va='center',
                      fontsize=12, fontweight='bold',
                      bbox=dict(boxstyle='round', facecolor='wheat', alpha=0.5))
        ax_bottom.axis('off')

        if save_path is None:
            save_path = self.output_dir / "dashboard_kpis.png"

        plt.savefig(save_path, dpi=300, bbox_inches='tight')
        plt.close()

        logger.info(f"KPI dashboard chart saved: {save_path}")
        return str(save_path)

    # ========================================================================
    # POWERPOINT GENERATION
    # ========================================================================

    def generate_powerpoint_deck(
        self,
        output_path: Optional[str] = None
    ) -> str:
        """
        Generate PowerPoint executive deck.

        Args:
            output_path: Path to save PowerPoint file

        Returns:
            Path to saved file
        """
        logger.info("Generating PowerPoint deck...")

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Slide 1: Title Slide
        self._add_title_slide(prs)

        # Slide 2: Executive Summary
        self._add_executive_summary_slide(prs)

        # Slide 3: KPI Dashboard
        self._add_kpi_slide(prs)

        # Slide 4: Time Trends
        self._add_trends_slide(prs)

        # Slide 5: Facility Type Comparison
        self._add_facility_comparison_slide(prs)

        # Slide 6: Capacity Analysis
        self._add_capacity_slide(prs)

        # Slide 7: Recommendations
        self._add_recommendations_slide(prs)

        if output_path is None:
            timestamp = datetime.now().strftime("%Y%m%d")
            output_path = self.output_dir / f"VA_COVID_Executive_Brief_{timestamp}.pptx"

        prs.save(output_path)
        logger.info(f"PowerPoint deck saved: {output_path}")
        return str(output_path)

    def _add_title_slide(self, prs: Presentation) -> None:
        """Add title slide to presentation."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Add title
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = "VA COVID-19 Response"

        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(44)
        p.font.bold = True

        # Add subtitle
        top = Inches(3.8)
        textbox2 = slide.shapes.add_textbox(left, top, width, height)
        text_frame2 = textbox2.text_frame
        text_frame2.text = f"Executive Dashboard\n{self.summary['kpis']['report_date']}"

        p2 = text_frame2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(24)

    def _add_executive_summary_slide(self, prs: Presentation) -> None:
        """Add executive summary slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content

        title = slide.shapes.title
        title.text = "Executive Summary"

        kpis = self.summary['kpis']

        # Content
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        summary_items = [
            f"• Facilities Reporting: {kpis['facilities_reporting']} facilities",
            f"• Total COVID-19 Positive: {kpis['total_covid_positive']:,} cases",
            f"• Hospitalized Patients: {kpis['total_covid_hospitalized']:,}",
            f"• ICU Patients: {kpis['total_covid_icu']:,}",
            f"• Patients on Ventilators: {kpis['total_covid_ventilator']:,}",
            f"• Deaths: {kpis['total_covid_deaths']:,}",
            f"• Average Positivity Rate: {kpis['avg_positivity_rate']:.1%}",
            f"• Average Occupancy Rate: {kpis['avg_occupancy_rate']:.1%}",
            f"• Facilities with Staff Shortages: {kpis['facilities_staff_shortage']}",
            f"• Facilities with PPE Critical: {kpis['facilities_ppe_critical']}",
        ]

        for item in summary_items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.space_before = Pt(6)

    def _add_kpi_slide(self, prs: Presentation) -> None:
        """Add KPI dashboard slide with chart."""
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only

        title = slide.shapes.title
        title.text = "Key Performance Indicators"

        # Generate and insert KPI chart
        chart_path = self.create_kpi_dashboard_chart()

        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5.5)

        slide.shapes.add_picture(chart_path, left, top, width=width, height=height)

    def _add_trends_slide(self, prs: Presentation) -> None:
        """Add trends slide with time series charts."""
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        title = slide.shapes.title
        title.text = "COVID-19 Trends Over Time"

        # Generate time series chart
        chart_path = self.create_time_series_chart(
            metric='covid_positive',
            title='COVID-19 Positive Cases - 7-Day Trend',
            ylabel='Cases'
        )

        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5.5)

        slide.shapes.add_picture(chart_path, left, top, width=width, height=height)

    def _add_facility_comparison_slide(self, prs: Presentation) -> None:
        """Add facility type comparison slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        title = slide.shapes.title
        title.text = "Facility Type Comparison"

        chart_path = self.create_facility_type_comparison_chart()

        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5.5)

        slide.shapes.add_picture(chart_path, left, top, width=width, height=height)

    def _add_capacity_slide(self, prs: Presentation) -> None:
        """Add capacity analysis slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title = slide.shapes.title
        title.text = "Capacity Analysis"

        capacity = self.summary['capacity_analysis']

        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        items = [
            "Bed Capacity:",
            f"  • Total Beds: {capacity['total_beds_system']:,}",
            f"  • Occupied: {capacity['total_occupied_beds']:,}",
            f"  • Available: {capacity['total_available_beds']:,}",
            f"  • Average Occupancy: {capacity['avg_occupancy_rate']:.1%}",
            "",
            "Resource Constraints:",
            f"  • High Occupancy Facilities: {capacity['facilities_high_occupancy']}",
            f"  • Staff Shortages: {capacity['facilities_with_staff_shortage']} facilities",
            f"  • PPE Critical: {capacity['facilities_with_ppe_critical']} facilities",
        ]

        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.space_before = Pt(4)

    def _add_recommendations_slide(self, prs: Presentation) -> None:
        """Add recommendations slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title = slide.shapes.title
        title.text = "Recommendations"

        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        capacity = self.summary['capacity_analysis']
        kpis = self.summary['kpis']

        recommendations = []

        # Dynamic recommendations based on data
        if capacity['facilities_high_occupancy'] > 10:
            recommendations.append(
                "• URGENT: Deploy additional staff to high-occupancy facilities"
            )

        if kpis['facilities_ppe_critical'] > 5:
            recommendations.append(
                "• Expedite PPE supply shipments to facilities with critical shortages"
            )

        if kpis['avg_positivity_rate'] > 0.10:
            recommendations.append(
                "• Increase testing capacity - positivity rate exceeds 10% threshold"
            )

        if capacity['facilities_with_staff_shortage'] > 15:
            recommendations.append(
                "• Activate staffing contingency plans for affected facilities"
            )

        # Default recommendations
        recommendations.extend([
            "• Continue daily monitoring of capacity constraints",
            "• Maintain PPE inventory at 14+ day supply levels",
            "• Review staffing models for facilities with recurring shortages",
        ])

        for rec in recommendations:
            p = tf.add_paragraph()
            p.text = rec
            p.font.size = Pt(16)
            p.space_before = Pt(8)

    # ========================================================================
    # EXCEL WORKBOOK GENERATION
    # ========================================================================

    def generate_excel_workbook(
        self,
        output_path: Optional[str] = None
    ) -> str:
        """
        Generate Excel workbook with formatted data and charts.

        Args:
            output_path: Path to save Excel file

        Returns:
            Path to saved file
        """
        logger.info("Generating Excel workbook...")

        # Export data to Excel
        df = self.analyzer.df

        if output_path is None:
            timestamp = datetime.now().strftime("%Y%m%d")
            output_path = self.output_dir / f"VA_COVID_Data_{timestamp}.xlsx"

        # Create Excel writer
        with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
            # Sheet 1: Executive Summary
            summary_df = pd.DataFrame([self.summary['kpis']])
            summary_df.to_excel(writer, sheet_name='Executive Summary', index=False)

            # Sheet 2: Daily Data
            df.to_excel(writer, sheet_name='Daily Data', index=False)

            # Sheet 3: Facility Type Comparison
            comparison = self.analyzer.compare_by_facility_type()
            comparison.to_excel(writer, sheet_name='Facility Comparison', index=False)

            # Sheet 4: Capacity Analysis
            capacity = pd.DataFrame([self.summary['capacity_analysis']])
            capacity.to_excel(writer, sheet_name='Capacity Analysis', index=False)

            # Sheet 5: Time Trends
            trends = self.analyzer.analyze_time_trends('covid_positive')
            trends.to_excel(writer, sheet_name='COVID Trends', index=False)

        # Format workbook
        self._format_excel_workbook(output_path)

        logger.info(f"Excel workbook saved: {output_path}")
        return str(output_path)

    def _format_excel_workbook(self, file_path: str) -> None:
        """Apply formatting to Excel workbook."""
        wb = openpyxl.load_workbook(file_path)

        # Header style
        header_font = Font(bold=True, color="FFFFFF", size=11)
        header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
        header_alignment = Alignment(horizontal="center", vertical="center")

        # Apply header formatting to all sheets
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]

            # Format headers (first row)
            for cell in ws[1]:
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = header_alignment

            # Auto-adjust column widths
            for column in ws.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                ws.column_dimensions[column_letter].width = adjusted_width

        wb.save(file_path)

    # ========================================================================
    # COMPLETE DASHBOARD GENERATION
    # ========================================================================

    def generate_all_reports(self) -> Dict[str, str]:
        """
        Generate all reports and dashboards.

        Returns:
            Dictionary with paths to generated files
        """
        logger.info("=" * 80)
        logger.info("Generating all reports and dashboards...")
        logger.info("=" * 80)

        outputs = {}

        # Generate PowerPoint deck
        logger.info("\n[1/3] Generating PowerPoint deck...")
        outputs['powerpoint'] = self.generate_powerpoint_deck()

        # Generate Excel workbook
        logger.info("\n[2/3] Generating Excel workbook...")
        outputs['excel'] = self.generate_excel_workbook()

        # Generate standalone charts
        logger.info("\n[3/3] Generating standalone charts...")
        outputs['kpi_chart'] = self.create_kpi_dashboard_chart()
        outputs['trends_chart'] = self.create_time_series_chart(
            'covid_positive',
            'COVID-19 Positive Cases Over Time',
            'Cases'
        )
        outputs['comparison_chart'] = self.create_facility_type_comparison_chart()

        logger.info("\n" + "=" * 80)
        logger.info("All reports generated successfully!")
        logger.info("=" * 80)
        for report_type, path in outputs.items():
            logger.info(f"  {report_type}: {path}")
        logger.info("=" * 80)

        return outputs


# ============================================================================
# EXAMPLE USAGE
# ============================================================================

if __name__ == "__main__":
    from etl_pipeline import FacilityDataETL
    from analysis import CovidFacilityAnalyzer

    # Run ETL
    etl = FacilityDataETL()
    etl.run_pipeline(
        ehr_file="data/ehr_data.csv",
        staffing_file="data/staffing_roster.csv",
        ppe_file="data/ppe_inventory.csv",
        output_file="output/merged_data.csv"
    )

    # Run analysis
    analyzer = CovidFacilityAnalyzer(etl.merged_snapshots)

    # Generate reports
    generator = DashboardGenerator(analyzer, output_dir="output/reports")
    outputs = generator.generate_all_reports()

    print("\n✅ All reports generated successfully!")
    print("\nGenerated files:")
    for report_type, path in outputs.items():
        print(f"  - {report_type}: {path}")
